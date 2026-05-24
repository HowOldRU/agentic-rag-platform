"""
PPT 文档提取工具 - 使用 python-pptx 库
"""
import os
import uuid
import zipfile
import logging
from typing import Dict, Any
from abc import ABC
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)


class PPTExtractionTool(ABC):
    """PPT 提取工具，使用 python-pptx 库解析 PPT 文件"""
    name: str = "ppt_extractor"
    description: str = "PPT 文件提取工具，支持提取 PPT 中的文本内容，支持 .pptx 格式"

    def __init__(self):
        pass

    def is_ppt_file(self, file_name: str) -> bool:
        """检查是否为 PPT 文件"""
        if not file_name:
            return False
        return file_name.lower().endswith((".pptx", ".ppt"))

    def extract_text_from_slide(self, slide) -> str:
        """从幻灯片中提取文本"""
        text_runs = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                # 处理表格
                for row in shape.table.rows:
                    for cell in row.cells:
                        text_runs.append(cell.text)
        return '\n'.join(text_runs)

    def validate_ppt_file(self, file_path: str) -> Dict[str, Any]:
        """验证 PPT 文件是否有效"""
        try:
            # 检查文件大小
            file_size = os.path.getsize(file_path)
            if file_size == 0:
                return {
                    "valid": False,
                    "error": "文件为空"
                }

            # 对于 pptx 文件，检查是否为有效的 zip 文件
            if file_path.lower().endswith('.pptx'):
                try:
                    with zipfile.ZipFile(file_path, 'r') as zip_ref:
                        # 尝试读取 zip 文件内容来验证其有效性
                        zip_ref.testzip()
                except zipfile.BadZipFile:
                    return {
                        "valid": False,
                        "error": "文件不是有效的 ZIP 文件，可能已损坏"
                    }
                except Exception as e:
                    return {
                        "valid": False,
                        "error": f"ZIP 文件验证失败: {str(e)}"
                    }

            return {
                "valid": True,
                "error": None
            }
        except Exception as e:
            return {
                "valid": False,
                "error": f"文件验证过程中发生错误: {str(e)}"
            }

    def run(self, **kwargs) -> Dict[str, Any]:
        """执行 PPT 提取，使用 python-pptx 库解析 PPT 文件"""
        file_content = kwargs.get("file_content")
        file_name = kwargs.get("file_name", "")
        temp_path = None

        try:
            # 验证文件类型
            if not self.is_ppt_file(file_name):
                return {
                    "success": False,
                    "error": f"不支持的文件类型，仅支持 PPTX 文件。当前文件: {file_name}"
                }

            # 检查文件内容是否为空
            if not file_content:
                return {
                    "success": False,
                    "error": "文件内容为空"
                }

            # 生成唯一的临时文件名以避免冲突
            file_ext = file_name.lower().split('.')[-1]
            temp_path = f"temp_{uuid.uuid4().hex}.{file_ext}"

            # 临时保存文件内容
            with open(temp_path, "wb") as f:
                f.write(file_content)

            # 检查文件是否创建成功且不为空
            if not os.path.exists(temp_path) or os.path.getsize(temp_path) == 0:
                return {
                    "success": False,
                    "error": "临时文件创建失败或文件为空"
                }

            # 验证 PPT 文件是否有效
            validation_result = self.validate_ppt_file(temp_path)
            if not validation_result["valid"]:
                return {
                    "success": False,
                    "error": validation_result["error"]
                }

            # 使用 python-pptx 库解析 PPT 文件
            prs = Presentation(temp_path)

            # 提取幻灯片内容
            serialized_documents = []
            for slide_idx, slide in enumerate(prs.slides, 1):
                # 提取幻灯片文本
                slide_text = self.extract_text_from_slide(slide)

                # 创建元数据
                metadata = {
                    "slide_number": slide_idx,
                    "file_type": "pptx",
                    "file_name": file_name,
                    "slide_count": len(prs.slides)
                }

                # 创建文档对象
                serialized_doc = {
                    "page_content": slide_text,
                    "metadata": metadata
                }
                serialized_documents.append(serialized_doc)

            result = {
                "success": True,
                "documents": serialized_documents,
                "slide_count": len(serialized_documents)
            }

            return result

        except zipfile.BadZipFile as e:
            error_msg = f"PPT 文件处理失败: 文件已损坏或不是有效的 PPTX 文件"
            logger.error(error_msg)
            return {
                "success": False,
                "error": error_msg,
                "documents": []
            }
        except Exception as e:
            error_msg = f"PPT 文件处理失败: {str(e)}"
            logger.error(error_msg)
            return {
                "success": False,
                "error": error_msg,
                "documents": []
            }
        finally:
            # 确保清理临时文件
            try:
                if temp_path and os.path.exists(temp_path):
                    os.remove(temp_path)
            except:
                pass
