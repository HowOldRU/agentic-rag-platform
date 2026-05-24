"""
PPT 渲染引擎

纯渲染模块，负责：
1. 根据主题和内容规划生成幻灯片结构（build_slide_plan）
2. 使用 python-pptx 渲染为 .pptx 字节流（render_presentation）

不依赖 BaseTool，可独立调用和测试。
"""

from io import BytesIO
from uuid import uuid4

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

from .ppt_themes import (
    THEMES,
    DEFAULT_THEME,
    SLIDE_WIDTH_16_9,
    SLIDE_HEIGHT_16_9,
    MARGIN_LEFT,
    MARGIN_TOP,
    CONTENT_WIDTH,
)


# ========== 幻灯片结构生成 ==========

def build_slide_plan(
    topic: str,
    slide_count: int = 5,
    slides: list[dict] | None = None,
) -> dict:
    """
    根据主题和参数生成幻灯片结构计划。

    Args:
        topic: 演示文稿主题
        slide_count: 内容页数量（3-15），不含自动添加的标题页/目录/结尾
        slides: LLM 提供的结构化幻灯片内容列表。
            每项 {"title": "...", "bullets": ["...", "..."]}。
            如果提供，将直接使用这些内容，忽略 slide_count。

    Returns:
        结构化幻灯片计划字典
    """
    slides_out = []

    # 1. 标题页
    import datetime
    today = datetime.date.today().strftime("%Y年%m月%d日")
    slides_out.append({
        "layout": "title",
        "title": topic,
        "subtitle": today,
    })

    # 2. 内容页（来自 LLM 或自动生成）
    if slides and isinstance(slides, list) and len(slides) > 0:
        # LLM 提供了具体内容，直接使用
        content_slides = slides
    else:
        # 无 LLM 内容时，用通用模板兜底
        content_slides = _fallback_sections(topic, slide_count)

    # 目录页（从内容页标题生成）
    slides_out.append({
        "layout": "toc",
        "title": "目录",
        "items": [s.get("title", "") for s in content_slides],
    })

    # 内容页
    for i, slide_data in enumerate(content_slides):
        # 每 4 页插入一个章节分隔页
        if i > 0 and i % 4 == 0 and len(content_slides) > 5:
            slides_out.append({
                "layout": "section",
                "title": slide_data.get("title", ""),
            })
        slides_out.append({
            "layout": "content",
            "title": slide_data.get("title", ""),
            "bullets": slide_data.get("bullets", []),
        })

    # 3. 结尾页
    slides_out.append({
        "layout": "ending",
        "title": "谢谢",
        "subtitle": f"关于「{topic}」的分享到此结束，欢迎交流讨论",
    })

    return {
        "title": topic,
        "slide_count": len(slides_out),
        "slides": slides_out,
    }


def _fallback_sections(topic: str, count: int) -> list[dict]:
    """LLM 未提供内容时的兜底模板"""
    count = max(3, min(15, count))
    templates = [
        {"title": f"{topic}概述", "bullets": [f"{topic}的定义与基本概念", "发展背景与意义", "核心价值"]},
        {"title": "现状分析", "bullets": ["当前发展态势", "关键数据指标", "主要特征"]},
        {"title": "核心内容", "bullets": ["关键组成部分", "基本原则与方法", "重要成果"]},
        {"title": "实施路径", "bullets": ["总体策略与框架", "关键步骤", "预期成效"]},
        {"title": "总结展望", "bullets": ["核心观点回顾", "未来发展方向", "行动建议"]},
    ]
    return templates[:count]


# ========== PPT 渲染 ==========

def render_presentation(plan: dict, theme_name: str = DEFAULT_THEME) -> bytes:
    """
    根据幻灯片计划渲染为 .pptx 字节流。

    Args:
        plan: build_slide_plan() 返回的结构计划
        theme_name: 主题名称

    Returns:
        .pptx 文件的字节数据
    """
    theme = THEMES.get(theme_name, THEMES[DEFAULT_THEME])

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH_16_9
    prs.slide_height = SLIDE_HEIGHT_16_9

    # 获取空白布局
    blank_layout = prs.slide_layouts[6]

    for slide_data in plan.get("slides", []):
        layout_type = slide_data.get("layout", "content")
        slide = prs.slides.add_slide(blank_layout)

        if layout_type == "title":
            _render_title_slide(slide, slide_data, theme)
        elif layout_type == "toc":
            _render_toc_slide(slide, slide_data, theme)
        elif layout_type == "content":
            _render_content_slide(slide, slide_data, theme)
        elif layout_type == "section":
            _render_section_slide(slide, slide_data, theme)
        elif layout_type == "ending":
            _render_ending_slide(slide, slide_data, theme)
        else:
            _render_content_slide(slide, slide_data, theme)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ========== 各布局渲染函数 ==========

def _add_gradient_bg(
    slide,
    color_start,
    color_end,
) -> None:
    """添加全页渐变背景（矩形形状）"""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        0, 0,
        SLIDE_WIDTH_16_9, SLIDE_HEIGHT_16_9,
    )
    shape.line.fill.background()  # 无边框

    # 设置渐变填充
    fill = shape.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = color_start
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = color_end
    fill.gradient_stops[1].position = 1.0

    # 移到最底层
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)


def _add_textbox(
    slide,
    text: str,
    left: int,
    top: int,
    width: int,
    height: int,
    font_size=Pt(18),
    font_color=None,
    bold: bool = False,
    alignment=PP_ALIGN.LEFT,
    font_name: str = "微软雅黑",
) -> None:
    """添加文本框"""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if font_color:
        p.font.color.rgb = font_color


def _add_header_bar(
    slide,
    title_text: str,
    theme: dict,
) -> None:
    """添加顶部标题栏（彩色条带 + 白色标题文字）"""
    bar_height = Inches(0.9)
    # 背景条
    bar = slide.shapes.add_shape(
        1, 0, 0,
        SLIDE_WIDTH_16_9, bar_height,
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = theme["header_bg"]
    bar.line.fill.background()

    # 标题文字
    _add_textbox(
        slide,
        title_text,
        MARGIN_LEFT, Inches(0.12),
        CONTENT_WIDTH, bar_height - Inches(0.2),
        font_size=theme["title_size"],
        font_color=theme["header_text_color"],
        bold=True,
        alignment=PP_ALIGN.LEFT,
    )


def _add_bullets(
    slide,
    bullets: list[str],
    top: int,
    theme: dict,
) -> None:
    """添加项目符号列表"""
    if not bullets:
        return

    body_size = theme["body_size"]
    line_height = int(body_size * 2.2)
    total_height = line_height * len(bullets) + Inches(0.3)

    txbox = slide.shapes.add_textbox(
        MARGIN_LEFT + Inches(0.3), top,
        CONTENT_WIDTH - Inches(0.3), total_height,
    )
    tf = txbox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = body_size
        p.font.color.rgb = theme["body_color"]
        p.font.name = "微软雅黑"
        p.space_after = Pt(6)
        # 添加圆点项目符号
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        buChar = pPr.makeelement(qn("a:buChar"), {"char": "•"})
        # 移除已有符号设置
        for child in list(pPr):
            if child.tag.endswith("buChar") or child.tag.endswith("buNone"):
                pPr.remove(child)
        pPr.append(buChar)


def _render_title_slide(slide, data: dict, theme: dict) -> None:
    """标题页：渐变背景 + 居中大标题 + 副标题"""
    _add_gradient_bg(slide, theme["bg_start"], theme["bg_end"])

    # 装饰线
    line_width = Inches(2.0)
    line_left = (SLIDE_WIDTH_16_9 - line_width) // 2
    line = slide.shapes.add_shape(
        1, line_left, Inches(3.4),
        line_width, Pt(3),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = theme["primary"]
    line.line.fill.background()

    # 标题
    title_text = data.get("title", "")
    _add_textbox(
        slide,
        title_text,
        Inches(1), Inches(2.0),
        SLIDE_WIDTH_16_9 - Inches(2), Inches(1.5),
        font_size=Pt(44),
        font_color=theme["title_color"],
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    # 副标题
    subtitle = data.get("subtitle", "")
    _add_textbox(
        slide,
        subtitle,
        Inches(1.5), Inches(3.8),
        SLIDE_WIDTH_16_9 - Inches(3), Inches(1.0),
        font_size=theme["subtitle_size"],
        font_color=theme["secondary"],
        bold=False,
        alignment=PP_ALIGN.CENTER,
    )


def _render_toc_slide(slide, data: dict, theme: dict) -> None:
    """目录页：标题栏 + 编号列表"""
    _add_gradient_bg(slide, theme["bg_start"], theme["bg_end"])
    _add_header_bar(slide, data.get("title", "目录"), theme)

    items = data.get("items", [])
    if not items:
        return

    # 目录项 — 两列布局
    col_count = 2
    col_width = (CONTENT_WIDTH - Inches(0.6)) // col_count
    start_top = Inches(1.4)

    for i, item in enumerate(items):
        col = i % col_count
        row = i // col_count
        left = MARGIN_LEFT + col * (col_width + Inches(0.6))
        top = start_top + row * Inches(0.7)

        # 编号圆圈（小方块模拟）
        num_size = Inches(0.35)
        num_shape = slide.shapes.add_shape(
            1, left, top + Pt(3), num_size, num_size,
        )
        num_shape.fill.solid()
        num_shape.fill.fore_color.rgb = theme["primary"]
        num_shape.line.fill.background()
        ntf = num_shape.text_frame
        ntf.paragraphs[0].text = str(i + 1)
        ntf.paragraphs[0].font.size = Pt(14)
        ntf.paragraphs[0].font.color.rgb = theme["header_text_color"]
        ntf.paragraphs[0].font.bold = True
        ntf.paragraphs[0].alignment = PP_ALIGN.CENTER
        ntf.word_wrap = False

        # 目录文字
        _add_textbox(
            slide,
            item,
            left + num_size + Inches(0.15), top,
            col_width - num_size - Inches(0.15), Inches(0.5),
            font_size=Pt(18),
            font_color=theme["body_color"],
            bold=False,
        )


def _render_content_slide(slide, data: dict, theme: dict) -> None:
    """内容页：标题栏 + 项目符号列表"""
    _add_gradient_bg(slide, theme["bg_start"], theme["bg_end"])
    _add_header_bar(slide, data.get("title", ""), theme)

    bullets = data.get("bullets", [])
    _add_bullets(slide, bullets, Inches(1.4), theme)


def _render_section_slide(slide, data: dict, theme: dict) -> None:
    """章节分隔页：纯色背景 + 居中标题"""
    # 全屏主色背景
    bg = slide.shapes.add_shape(
        1, 0, 0,
        SLIDE_WIDTH_16_9, SLIDE_HEIGHT_16_9,
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme["primary"]
    bg.line.fill.background()

    # 居中标题
    _add_textbox(
        slide,
        data.get("title", ""),
        Inches(1), (SLIDE_HEIGHT_16_9 - Inches(1.5)) // 2,
        SLIDE_WIDTH_16_9 - Inches(2), Inches(1.5),
        font_size=Pt(40),
        font_color=theme["header_text_color"],
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )


def _render_ending_slide(slide, data: dict, theme: dict) -> None:
    """结尾页：渐变背景 + 居中感谢语"""
    _add_gradient_bg(slide, theme["bg_start"], theme["bg_end"])

    # 主标题
    _add_textbox(
        slide,
        data.get("title", "谢谢"),
        Inches(1), Inches(2.4),
        SLIDE_WIDTH_16_9 - Inches(2), Inches(1.2),
        font_size=Pt(48),
        font_color=theme["title_color"],
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    # 副标题
    subtitle = data.get("subtitle", "")
    if subtitle:
        _add_textbox(
            slide,
            subtitle,
            Inches(1.5), Inches(3.8),
            SLIDE_WIDTH_16_9 - Inches(3), Inches(0.8),
            font_size=Pt(18),
            font_color=theme["secondary"],
            bold=False,
            alignment=PP_ALIGN.CENTER,
        )
